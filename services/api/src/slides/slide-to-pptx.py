"""
Convert an HTML slide presentation to an editable PPTX via python-pptx + lxml.

Usage:  cat slides.html | python slide-to-pptx.py > output.pptx

Each .slide / .page / section element becomes one PowerPoint slide.
Text is extracted and placed in real text boxes — fully editable in PowerPoint.
"""
import sys
import re
import io
from lxml import html as lhtml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.55)


# ── Colour helpers ────────────────────────────────────────────────────────────

def _parse_hex(h: str):
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) == 6:
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            pass
    return None


NAMED_COLORS = {
    "white": "ffffff", "black": "000000", "navy": "001f3f",
    "darkblue": "00008b", "dark": "0b0f17", "transparent": None,
}

CSS_RGB_RE = re.compile(
    r"rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", re.I
)


def parse_color(css: str):
    if not css:
        return None
    css = css.strip().lower()
    if css in ("none", "transparent", "inherit", "initial", "unset"):
        return None
    if css.startswith("#"):
        return _parse_hex(css)
    m = CSS_RGB_RE.match(css)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    for name, hexval in NAMED_COLORS.items():
        if name in css and hexval:
            return _parse_hex(hexval)
    return None


def luminance(rgb: RGBColor) -> float:
    r, g, b = rgb[0] / 255, rgb[1] / 255, rgb[2] / 255
    return 0.299 * r + 0.587 * g + 0.114 * b


def text_color_for(bg: RGBColor) -> RGBColor:
    return RGBColor(255, 255, 255) if luminance(bg) < 0.45 else RGBColor(20, 20, 20)


DEFAULT_BG = RGBColor(11, 15, 23)   # matches prodeck dark theme (#0b0f17)


# ── Slide content extraction ──────────────────────────────────────────────────

BG_RE = re.compile(r"background(?:-color)?\s*:\s*([^;]+)", re.I)


def extract_bg(el) -> RGBColor:
    style = el.get("style", "")
    m = BG_RE.search(style)
    if m:
        c = parse_color(m.group(1).strip())
        if c:
            return c
    # Check class-derived hints (e.g. bg-dark, bg-white)
    cls = " ".join(el.get("class", "").split())
    if "light" in cls or "white" in cls:
        return RGBColor(255, 255, 255)
    return DEFAULT_BG


BLOCK_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "dt", "dd"}
SKIP_TAGS  = {"script", "style", "svg", "noscript", "nav"}


def _text_of(el) -> str:
    """Recursively collect visible text, skipping style/script subtrees."""
    tag = (el.tag or "").lower() if isinstance(el.tag, str) else ""
    if tag in SKIP_TAGS:
        return ""
    parts = []
    if el.text:
        parts.append(el.text.strip())
    for child in el:
        parts.append(_text_of(child))
        if child.tail:
            parts.append(child.tail.strip())
    return " ".join(p for p in parts if p)


def extract_blocks(slide_el) -> list:
    """
    Walk the slide element and return a list of (kind, text) tuples.
    kind is one of: 'title', 'subtitle', 'body'
    """
    seen: set = set()
    blocks: list = []

    def visit(el, depth: int = 0) -> None:
        tag = (el.tag or "").lower() if isinstance(el.tag, str) else ""
        if tag in SKIP_TAGS:
            return

        if tag in BLOCK_TAGS:
            text = _text_of(el).strip()
            if text and text not in seen and len(text) < 800:
                seen.add(text)
                if tag == "h1":
                    kind = "title"
                elif tag in ("h2", "h3"):
                    kind = "subtitle"
                else:
                    kind = "body"
                blocks.append((kind, text))
        else:
            for child in el:
                visit(child, depth + 1)

    visit(slide_el)
    return blocks[:30]


# ── PPTX construction ─────────────────────────────────────────────────────────

def build_slide(prs: Presentation, bg: RGBColor, blocks: list) -> None:
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)

    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = bg

    fg = text_color_for(bg)
    content_w = SLIDE_W - 2 * MARGIN

    y = MARGIN

    for kind, text in blocks:
        if kind == "title":
            font_pt, line_h, bold = 40, Inches(0.75), True
        elif kind == "subtitle":
            font_pt, line_h, bold = 26, Inches(0.5),  True
        else:
            font_pt, line_h, bold = 18, Inches(0.38), False

        # Stop adding blocks if we'd overflow the slide
        if y + line_h > SLIDE_H - MARGIN:
            break

        txBox = slide.shapes.add_textbox(MARGIN, y, content_w, line_h)
        tf    = txBox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT

        run            = para.add_run()
        run.text       = text
        run.font.size  = Pt(font_pt)
        run.font.bold  = bold
        run.font.color.rgb = fg

        y += line_h + Inches(0.08)


# ── Entry point ───────────────────────────────────────────────────────────────

SLIDE_CLASS_TOKENS = {"slide", "page", "slide-page"}


def find_slide_elements(tree):
    # 1. data-slide / data-idx attributes
    els = tree.xpath("//*[@data-slide or @data-idx]")
    if els:
        return els

    # 2. standalone class tokens
    els = [
        el for el in tree.iter()
        if isinstance(el.tag, str)
        and SLIDE_CLASS_TOKENS & set(el.get("class", "").split())
    ]
    if els:
        return els

    # 3. <section> elements
    els = tree.xpath("//section")
    if els:
        return els

    # 4. fallback — body itself
    body = tree.xpath("//body")
    return body if body else [tree]


def main() -> None:
    html_bytes = sys.stdin.buffer.read()
    html_str   = html_bytes.decode("utf-8")

    tree      = lhtml.fromstring(html_str)
    slide_els = find_slide_elements(tree)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for el in slide_els:
        bg     = extract_bg(el)
        blocks = extract_blocks(el)
        build_slide(prs, bg, blocks)

    buf = io.BytesIO()
    prs.save(buf)
    sys.stdout.buffer.write(buf.getvalue())


if __name__ == "__main__":
    main()
